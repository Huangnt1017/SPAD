"""生成经过本地渲染迭代的 SPAD 可编辑网络结构 PPT。

CLI:
    D:\\Anaconda3\\envs\\torchnew\\python.exe scripts\\generate_model_architecture_ppt.py --output-dir model\\ppt

所有模块、文字和连接线均为 PowerPoint 原生可编辑对象，只输出一个汇总版 PPTX。
"""

from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Optional, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


W, H = 13.333, 7.5
FONT = "Microsoft YaHei"

BG = RGBColor(246, 248, 252)
NAVY = RGBColor(18, 33, 58)
TEXT = RGBColor(37, 47, 64)
MUTED = RGBColor(105, 116, 135)
BORDER = RGBColor(218, 224, 234)
WHITE = RGBColor(255, 255, 255)

FEATURE = RGBColor(47, 107, 255)
FEATURE_BG = RGBColor(234, 240, 255)
PHYSICAL = RGBColor(0, 150, 150)
PHYSICAL_BG = RGBColor(229, 247, 247)
FUSION = RGBColor(235, 145, 23)
FUSION_BG = RGBColor(255, 245, 224)
RESIDUAL = RGBColor(112, 71, 190)
RESIDUAL_BG = RGBColor(242, 236, 253)
OUTPUT = RGBColor(31, 157, 85)
OUTPUT_BG = RGBColor(232, 247, 238)
ABLATE = RGBColor(221, 67, 75)
ABLATE_BG = RGBColor(253, 236, 238)
NEUTRAL_BG = RGBColor(241, 244, 249)


@dataclass
class ScriptConfig:
    output_dir: Path


def no_shadow(shape) -> None:
    """关闭主题继承阴影，保持扁平风格。"""
    try:
        shape.shadow.inherit = False
    except (AttributeError, ValueError):
        pass


def arrow_end(line) -> None:
    ln = line._get_or_add_ln()
    arrow = OxmlElement("a:tailEnd")
    arrow.set("type", "triangle")
    arrow.set("w", "sm")
    arrow.set("len", "sm")
    ln.append(arrow)


def add_text(slide, x, y, w, h, text, *, size=10, bold=False, color=TEXT,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    no_shadow(shape)
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return shape


def add_shape(slide, x, y, w, h, text, *, fill=WHITE, line=BORDER, accent=None,
              size=10, bold=False, font_color=TEXT, radius=0.12, dashed=False):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    no_shadow(shape)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line; shape.line.width = Pt(1.2)
    if dashed: shape.line.dash_style = 2
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(7 if accent else 4); tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = font_color
    if accent:
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y+0.08), Inches(0.055), Inches(max(0.1,h-0.16)))
        no_shadow(bar); bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    return shape


def add_panel(slide, x, y, w, h, title):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    no_shadow(panel); panel.fill.solid(); panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = BORDER; panel.line.width = Pt(1)
    add_text(slide, x+0.20, y+0.10, w-0.40, 0.30, title, size=11, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    return panel


def line(slide, x1, y1, x2, y2, *, color=MUTED, width=1.35, arrow=True, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    no_shadow(c); c.line.color.rgb = color; c.line.width = Pt(width)
    if dashed: c.line.dash_style = 2
    if arrow: arrow_end(c.line)
    return c


def route(slide, points, *, color=MUTED, width=1.35, dashed=False):
    """用多段水平/垂直直线构造稳定的正交数据流。"""
    for i, ((x1,y1),(x2,y2)) in enumerate(zip(points, points[1:])):
        line(slide, x1,y1,x2,y2, color=color, width=width, arrow=(i==len(points)-2), dashed=dashed)


def setup_slide(prs, title, subtitle, source):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = BG
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(W), Inches(0.68))
    no_shadow(header); header.fill.solid(); header.fill.fore_color.rgb = NAVY; header.line.fill.background()
    add_text(slide, 0.42,0.08,8.8,0.36,title,size=21,bold=True,color=WHITE,align=PP_ALIGN.LEFT)
    add_text(slide, 9.2,0.10,3.65,0.28,subtitle,size=9.5,color=RGBColor(218,226,240),align=PP_ALIGN.RIGHT)
    add_text(slide, 0.42,0.43,8.6,0.16,"Editable vector diagram · strict code mapping",size=7.7,color=RGBColor(170,186,211),align=PP_ALIGN.LEFT)
    add_text(slide,0.34,7.20,8.5,0.16,f"Source: {source}",size=7.3,color=MUTED,align=PP_ALIGN.LEFT)
    add_text(slide,10.45,7.20,2.5,0.16,"Snapshot: 2026-07-20",size=7.3,color=MUTED,align=PP_ALIGN.RIGHT)
    return slide


def pipeline(slide, block_text="4 Graph Blocks", agg_text="Multi-scale\nCat + Agg", head_text="MLP head / Centroid head\ncenter B×3"):
    y,h = 0.88,0.55
    items = [
        (0.38,1.15,"Input\nB×N×4",FEATURE_BG,FEATURE),
        (1.82,1.18,"Stem\n4→32→32",NEUTRAL_BG,NAVY),
        (3.28,1.65,block_text,FUSION_BG,FUSION),
        (5.23,1.55,agg_text,RESIDUAL_BG,RESIDUAL),
        (7.08,1.42,"Max + Avg\nGlobal Pool",FEATURE_BG,FEATURE),
        (8.80,1.40,"Classification\nlogits B×C",OUTPUT_BG,OUTPUT),
        (10.50,2.40,head_text,FUSION_BG,FUSION),
    ]
    for x,w,t,bg,c in items: add_shape(slide,x,y,w,h,t,fill=bg,line=c,accent=c,size=8.8,bold=True)
    for (x,w,*_), (nx,*__) in zip(items,items[1:]): line(slide,x+w,y+h/2,nx,y+h/2,color=NAVY,width=1.25)


def bezier(slide, start, control1, control2, end, *, color=MUTED, width=1.45, dashed=False, behind=None):
    """添加带显式控制点的可编辑三次 Bé塞尔曲线。"""
    points = [start, control1, control2, end]
    min_x = min(x for x, _ in points); max_x = max(x for x, _ in points)
    min_y = min(y for _, y in points); max_y = max(y for _, y in points)
    if max_x - min_x < 0.01: max_x = min_x + 0.01
    if max_y - min_y < 0.01: max_y = min_y + 0.01

    builder = slide.shapes.build_freeform(Inches(min_x), Inches(min_y))
    builder.add_line_segments([(Inches(max_x), Inches(max_y))])
    shape = builder.convert_to_shape()
    no_shadow(shape); shape.fill.background()
    shape.line.color.rgb = color; shape.line.width = Pt(width)
    if dashed: shape.line.dash_style = 2
    arrow_end(shape.line)

    path = shape._element.spPr.custGeom.pathLst[0]
    for child in list(path): path.remove(child)
    path.set("w", str(int(Inches(max_x - min_x))))
    path.set("h", str(int(Inches(max_y - min_y))))

    move = OxmlElement("a:moveTo")
    point = OxmlElement("a:pt")
    point.set("x", str(int(Inches(start[0] - min_x))))
    point.set("y", str(int(Inches(start[1] - min_y))))
    move.append(point); path.append(move)

    cubic = OxmlElement("a:cubicBezTo")
    for x, y in (control1, control2, end):
        point = OxmlElement("a:pt")
        point.set("x", str(int(Inches(x - min_x))))
        point.set("y", str(int(Inches(y - min_y))))
        cubic.append(point)
    path.append(cubic)

    if behind is not None:
        element = shape._element
        element.getparent().remove(element)
        behind._element.addnext(element)
    return shape


def pill(slide,x,y,w,text,color,bg):
    return add_shape(slide,x,y,w,0.32,text,fill=bg,line=color,size=7.8,bold=True)


def slide_edgeconv(prs):
    s=setup_slide(prs,"GraphResidual · Original","EdgeConv + local Q/K/V attention","model/graph_residual.py")
    pipeline(s,"Blocks ×4\n32→64→64→128→256")
    main_panel=add_panel(s,0.34,1.68,9.10,5.28,"GraphResidualBlock · local graph attention")
    add_panel(s,9.68,1.68,3.30,5.28,"What changes inside one block")

    # Feature row
    add_shape(s,0.62,2.36,0.78,0.52,"f\nB×Cin×N",fill=FEATURE_BG,line=FEATURE,accent=FEATURE,bold=True)
    add_shape(s,1.72,2.25,1.25,0.72,"Dynamic KNN\nfeature space",fill=FEATURE_BG,line=FEATURE,accent=FEATURE,size=9.2,bold=True)
    add_shape(s,3.28,2.25,1.38,0.72,"Graph feature\n[fj-fi, fi]",fill=FEATURE_BG,line=FEATURE,accent=FEATURE,size=9.2)
    add_shape(s,4.97,2.25,1.30,0.72,"EdgeConv_f\n→ Fk",fill=FEATURE_BG,line=FEATURE,accent=FEATURE,size=9.5,bold=True)
    add_shape(s,6.58,2.25,0.82,0.72,"V\nWv(Fk)",fill=RESIDUAL_BG,line=RESIDUAL,size=9.5,bold=True)

    # Physical row
    add_shape(s,0.62,4.42,0.78,0.52,"p\nB×4×N",fill=PHYSICAL_BG,line=PHYSICAL,accent=PHYSICAL,bold=True)
    add_shape(s,1.72,4.31,1.25,0.72,"Coordinate KNN\ncomputed once",fill=PHYSICAL_BG,line=PHYSICAL,accent=PHYSICAL,size=8.8,bold=True)
    add_shape(s,3.28,4.31,1.38,0.72,"Cached p_graph\n[pj-pi, pi]",fill=PHYSICAL_BG,line=PHYSICAL,accent=PHYSICAL,size=9.0)
    add_shape(s,4.97,4.31,1.30,0.72,"EdgeConv_p\n→ Pk",fill=PHYSICAL_BG,line=PHYSICAL,accent=PHYSICAL,size=9.5,bold=True)
    add_shape(s,6.58,4.31,0.82,0.72,"K\nWk(Pk)",fill=RESIDUAL_BG,line=RESIDUAL,size=9.5,bold=True)

    # Q and attention
    add_shape(s,3.72,3.30,1.30,0.58,"Q = Wq([f,p])",fill=RESIDUAL_BG,line=RESIDUAL,accent=RESIDUAL,size=9.5,bold=True)
    add_shape(s,7.70,2.94,1.34,1.10,"Local attention\nsoftmax over k\nΣ weight·V",fill=RESIDUAL_BG,line=RESIDUAL,accent=RESIDUAL,size=9.4,bold=True)
    add_shape(s,7.62,4.55,1.42,0.52,"Output mapping",fill=NEUTRAL_BG,line=NAVY,accent=NAVY,size=9.2)
    add_shape(s,5.00,5.48,1.22,0.54,"coord_gate(p)",fill=FUSION_BG,line=FUSION,accent=FUSION,size=9.0)
    add_shape(s,6.48,5.48,1.22,0.54,"coord_res(p)",fill=FUSION_BG,line=FUSION,accent=FUSION,size=9.0)
    add_shape(s,7.96,5.43,1.08,0.64,"Gated fusion",fill=FUSION_BG,line=FUSION,accent=FUSION,size=9.2,bold=True)
    add_shape(s,7.96,6.22,1.08,0.48,"Act → f_out",fill=OUTPUT_BG,line=OUTPUT,accent=OUTPUT,size=9.2,bold=True)

    # Main row routes
    route(s,[(1.40,2.62),(1.72,2.62)],color=FEATURE); route(s,[(2.97,2.62),(3.28,2.62)],color=FEATURE)
    route(s,[(4.66,2.62),(4.97,2.62)],color=FEATURE); route(s,[(6.27,2.62),(6.58,2.62)],color=FEATURE)
    route(s,[(1.40,4.68),(1.72,4.68)],color=PHYSICAL); route(s,[(2.97,4.68),(3.28,4.68)],color=PHYSICAL)
    route(s,[(4.66,4.68),(4.97,4.68)],color=PHYSICAL); route(s,[(6.27,4.68),(6.58,4.68)],color=PHYSICAL)
    # Q inputs and attention
    bezier(s,(1.40,2.74),(2.10,2.78),(2.85,3.30),(3.72,3.52),color=RESIDUAL,behind=main_panel)
    bezier(s,(1.40,4.56),(2.10,4.48),(2.85,3.92),(3.72,3.72),color=RESIDUAL,behind=main_panel)
    bezier(s,(5.02,3.59),(5.82,3.22),(6.82,3.20),(7.70,3.49),color=RESIDUAL,behind=main_panel)
    bezier(s,(7.40,2.62),(7.56,2.70),(7.58,3.02),(7.70,3.20),color=RESIDUAL,behind=main_panel)
    bezier(s,(7.40,4.68),(7.56,4.58),(7.58,3.98),(7.70,3.78),color=RESIDUAL,behind=main_panel)
    route(s,[(8.37,4.04),(8.37,4.55)],color=NAVY)
    route(s,[(8.33,5.07),(8.33,5.43)],color=FUSION)
    bezier(s,(1.01,4.94),(1.55,5.70),(3.60,5.98),(5.00,5.75),color=FUSION,behind=main_panel)
    bezier(s,(1.01,4.94),(2.05,6.48),(5.15,6.50),(6.48,5.75),color=FUSION,behind=main_panel)
    bezier(s,(6.22,5.75),(6.68,5.20),(7.38,5.18),(7.96,5.58),color=FUSION,behind=main_panel)
    bezier(s,(7.70,5.75),(7.78,5.78),(7.88,5.84),(7.96,5.88),color=FUSION,behind=main_panel)
    route(s,[(8.50,6.07),(8.50,6.22)],color=OUTPUT)

    # Notes
    pill(s,9.95,2.16,1.02,"FEATURE",FEATURE,FEATURE_BG); pill(s,11.12,2.16,1.25,"GEOMETRY",PHYSICAL,PHYSICAL_BG)
    add_text(s,9.95,2.60,2.72,0.78,"Two EdgeConv streams\nFeature graph is dynamic; coordinate graph is cached.",size=9.2,align=PP_ALIGN.LEFT)
    add_shape(s,9.95,3.55,2.72,0.82,"Q: [f,p]\nK: physical edge features\nV: semantic edge features",fill=RESIDUAL_BG,line=RESIDUAL,accent=RESIDUAL,size=9.1)
    add_shape(s,9.95,4.62,2.72,0.82,"gate·mapped\n+ (1-gate)·coord_info",fill=FUSION_BG,line=FUSION,accent=FUSION,size=9.5,bold=True)
    add_shape(s,9.95,5.68,2.72,0.82,"No point downsampling\nN remains 1024",fill=OUTPUT_BG,line=OUTPUT,accent=OUTPUT,size=9.5)
    add_text(s,9.95,6.60,2.72,0.20,"1.058M params with centroid head",size=8.5,bold=True,color=NAVY)
    return s


def slide_gcn(prs):
    s=setup_slide(prs,"GraphResidual-GCN · Formal","双图 GraphSAGE · SE · 受控双残差","model/graph_res_GCN.py")
    pipeline(s,"Blocks ×4\n32→64→64→128→256")
    main_panel=add_panel(s,0.34,1.68,9.10,5.28,"GraphResidualBlockGCN · non-legacy forward")
    add_panel(s,9.68,1.68,3.30,5.28,"Frozen formal configuration")

    # Two graph streams
    add_shape(s,0.62,2.30,0.78,0.52,"f",fill=FEATURE_BG,line=FEATURE,accent=FEATURE,bold=True)
    add_shape(s,1.72,2.19,1.18,0.72,"Feature KNN\nper block",fill=FEATURE_BG,line=FEATURE,accent=FEATURE,size=9.2,bold=True)
    add_shape(s,3.18,2.19,1.22,0.72,"SAGEConv_f\nAGG=max",fill=FEATURE_BG,line=FEATURE,accent=FEATURE,size=9.3,bold=True)
    add_shape(s,4.70,2.19,1.05,0.72,"BN + Act",fill=FEATURE_BG,line=FEATURE,size=9.3)
    add_shape(s,6.05,2.13,1.20,0.84,"SE gate\nGAP→MLP→σ",fill=RESIDUAL_BG,line=RESIDUAL,accent=RESIDUAL,size=9.2,bold=True)

    add_shape(s,0.62,4.26,0.78,0.52,"p",fill=PHYSICAL_BG,line=PHYSICAL,accent=PHYSICAL,bold=True)
    add_shape(s,1.72,4.15,1.18,0.72,"Coordinate KNN\ncached",fill=PHYSICAL_BG,line=PHYSICAL,accent=PHYSICAL,size=9.0,bold=True)
    add_shape(s,3.18,4.15,1.22,0.72,"SAGEConv_p\n4→Cout",fill=PHYSICAL_BG,line=PHYSICAL,accent=PHYSICAL,size=9.3,bold=True)
    add_shape(s,4.70,4.15,1.05,0.72,"BN + Act",fill=PHYSICAL_BG,line=PHYSICAL,size=9.3)

    # Fusion and output
    add_shape(s,7.55,3.00,0.72,0.72,"Cat",fill=FUSION_BG,line=FUSION,accent=FUSION,size=9.5,bold=True)
    add_shape(s,8.50,3.00,0.66,0.72,"Fuse\n2C→C",fill=FUSION_BG,line=FUSION,size=8.8,bold=True)
    add_shape(s,7.80,4.26,1.36,0.64,"Σ residuals\n→ Act → f_out",fill=OUTPUT_BG,line=OUTPUT,accent=OUTPUT,size=9.2,bold=True)

    # Clean residual rail
    add_shape(s,1.82,5.55,1.58,0.60,"feature residual\nIdentity / projection",fill=RESIDUAL_BG,line=RESIDUAL,accent=RESIDUAL,size=8.8)
    add_shape(s,3.92,5.55,1.18,0.60,"coord_gate",fill=FUSION_BG,line=FUSION,accent=FUSION,size=8.9)
    add_shape(s,5.40,5.55,1.76,0.60,"coord_res + encoder",fill=FUSION_BG,line=FUSION,accent=FUSION,size=8.8)
    add_shape(s,7.47,5.49,1.28,0.72,"gate × sum\n× coord_scale",fill=FUSION_BG,line=FUSION,accent=FUSION,size=8.5,bold=True)
    add_text(s,0.88,6.46,7.00,0.22,"PyG message flow:  neighbor/source j  ─────────▶  center/target i",size=8.8,bold=True,color=NAVY,align=PP_ALIGN.LEFT)

    # Main streams
    for a,b,c in [
        ((1.40,2.56),(1.72,2.56),FEATURE),((2.90,2.56),(3.18,2.56),FEATURE),
        ((4.40,2.56),(4.70,2.56),FEATURE),((5.75,2.56),(6.05,2.56),RESIDUAL),
        ((1.40,4.52),(1.72,4.52),PHYSICAL),((2.90,4.52),(3.18,4.52),PHYSICAL),
        ((4.40,4.52),(4.70,4.52),PHYSICAL),
    ]: route(s,[a,b],color=c)
    route(s,[(7.25,2.55),(7.38,2.55),(7.38,3.24),(7.55,3.24)],color=RESIDUAL)
    route(s,[(5.75,4.52),(7.38,4.52),(7.38,3.48),(7.55,3.48)],color=PHYSICAL)
    route(s,[(8.27,3.36),(8.50,3.36)],color=FUSION)
    route(s,[(8.83,3.72),(8.83,4.26)],color=OUTPUT)

    # Feature residual: below all modules, no diagonal crossing
    bezier(s,(1.01,2.82),(0.92,3.85),(1.05,5.45),(1.82,5.85),color=RESIDUAL,behind=main_panel)
    route(s,[(3.40,5.85),(3.55,5.85),(3.55,5.12),(7.55,5.12),(7.55,4.58),(7.80,4.58)],color=RESIDUAL)

    # Coordinate residual curves use separate upper/lower lanes.
    bezier(s,(1.01,4.78),(1.75,5.08),(3.25,5.08),(4.51,5.55),color=FUSION,behind=main_panel)
    bezier(s,(1.01,4.78),(2.15,6.45),(4.65,6.45),(6.28,5.55),color=FUSION,behind=main_panel)
    bezier(s,(5.10,5.85),(5.72,5.24),(6.82,5.20),(7.47,5.68),color=FUSION,behind=main_panel)
    bezier(s,(7.16,5.85),(7.26,5.88),(7.36,5.94),(7.47,5.98),color=FUSION,behind=main_panel)
    bezier(s,(8.11,5.49),(8.52,5.30),(8.62,5.04),(8.50,4.90),color=FUSION,behind=main_panel)

    # Right cards
    add_shape(s,9.95,2.10,2.72,0.72,"k=20 · AGG=max\nexclude_self=true",fill=FEATURE_BG,line=FEATURE,accent=FEATURE,size=9.4,bold=True)
    add_shape(s,9.95,3.03,2.72,0.72,"feature residual=true\ncoord_scale_init=0.1",fill=RESIDUAL_BG,line=RESIDUAL,accent=RESIDUAL,size=9.3)
    add_shape(s,9.95,3.96,2.72,0.72,"gradient checkpoint=true\nlegacy mode=false",fill=NEUTRAL_BG,line=NAVY,accent=NAVY,size=9.3)
    add_shape(s,9.95,4.89,2.72,0.82,"MLP head: pooled→3\nCentroid: point feats + xyz→3",fill=FUSION_BG,line=FUSION,accent=FUSION,size=9.1)
    add_shape(s,9.95,5.92,2.72,0.66,"1.101M · centroid head\n1.332M · MLP head",fill=OUTPUT_BG,line=OUTPUT,accent=OUTPUT,size=9.2,bold=True)
    return s

def slide_ablation(prs):
    s=setup_slide(prs,"GraphResidual-GCN · Controls","B0–B7 硬结构消融 · B8 参数匹配算子对照","model/graph_res_GCN_ablation.py · scripts/ablation_registry.py")
    pipeline(s,"Controlled Blocks ×4\nfull-width topology")
    main_panel=add_panel(s,0.34,1.68,8.80,5.28,"Controlled block · red dashed = hard removal · B8 changes operator only")
    add_panel(s,9.38,1.68,3.60,5.28,"Factor matrix · MLP · λobj=0")

    # Main streams
    add_shape(s,0.62,2.30,0.72,0.50,"f",fill=FEATURE_BG,line=FEATURE,accent=FEATURE,bold=True)
    add_shape(s,1.62,2.18,1.35,0.74,"Feature KNN\n+ local op_f",fill=FEATURE_BG,line=FEATURE,accent=FEATURE,size=9.2,bold=True)
    add_shape(s,3.28,2.18,1.22,0.74,"SE gate\noptional",fill=ABLATE_BG,line=ABLATE,accent=ABLATE,size=9.4,bold=True,dashed=True)

    add_shape(s,0.62,4.26,0.72,0.50,"p",fill=PHYSICAL_BG,line=PHYSICAL,accent=PHYSICAL,bold=True)
    add_shape(s,1.62,4.14,1.35,0.74,"Coordinate KNN\n+ local op_p",fill=ABLATE_BG,line=ABLATE,accent=ABLATE,size=8.9,bold=True,dashed=True)
    add_shape(s,3.28,4.14,1.22,0.74,"Physical branch\noptional",fill=ABLATE_BG,line=ABLATE,accent=ABLATE,size=9.0,bold=True,dashed=True)

    add_shape(s,4.92,3.06,0.88,0.74,"Cat /\nfeature only",fill=FUSION_BG,line=FUSION,accent=FUSION,size=8.8,bold=True)
    add_shape(s,6.10,3.06,0.76,0.74,"Fuse",fill=FUSION_BG,line=FUSION,accent=FUSION,size=9.4,bold=True)
    add_shape(s,7.18,3.06,0.84,0.74,"Σ + Act",fill=OUTPUT_BG,line=OUTPUT,accent=OUTPUT,size=9.2,bold=True)
    add_shape(s,8.28,3.06,0.58,0.74,"f_out",fill=OUTPUT_BG,line=OUTPUT,size=9.0,bold=True)

    # Residual rail with no diagonal line through labels
    add_shape(s,1.82,5.52,1.72,0.62,"feature residual\nexisting CLI toggle",fill=RESIDUAL_BG,line=RESIDUAL,accent=RESIDUAL,size=8.6)
    add_shape(s,4.02,5.42,2.22,0.82,"coordinate residual group\ncoord_gate + coord_res + encoder\noptional hard removal",fill=ABLATE_BG,line=ABLATE,accent=ABLATE,size=8.5,bold=True,dashed=True)
    add_shape(s,6.68,5.52,1.46,0.62,"coord_scale × delta",fill=ABLATE_BG,line=ABLATE,accent=ABLATE,size=8.6,dashed=True)

    route(s,[(1.34,2.55),(1.62,2.55)],color=FEATURE); route(s,[(2.97,2.55),(3.28,2.55)],color=FEATURE)
    route(s,[(4.50,2.55),(4.72,2.55),(4.72,3.28),(4.92,3.28)],color=ABLATE)
    route(s,[(1.34,4.51),(1.62,4.51)],color=PHYSICAL); route(s,[(2.97,4.51),(3.28,4.51)],color=ABLATE)
    route(s,[(4.50,4.51),(4.72,4.51),(4.72,3.58),(4.92,3.58)],color=ABLATE)
    route(s,[(5.80,3.43),(6.10,3.43)],color=FUSION); route(s,[(6.86,3.43),(7.18,3.43)],color=OUTPUT); route(s,[(8.02,3.43),(8.28,3.43)],color=OUTPUT)

    bezier(s,(0.98,2.80),(0.84,3.95),(0.98,5.52),(1.82,5.83),color=RESIDUAL,behind=main_panel)
    bezier(s,(3.54,5.83),(4.45,5.72),(5.78,4.10),(7.18,3.66),color=RESIDUAL,behind=main_panel)
    bezier(s,(0.98,4.76),(1.82,5.52),(3.92,5.08),(5.13,5.42),color=ABLATE,dashed=True,behind=main_panel)
    bezier(s,(6.24,5.83),(6.38,5.78),(6.54,5.78),(6.68,5.83),color=ABLATE,dashed=True,behind=main_panel)
    bezier(s,(8.14,5.83),(8.62,5.24),(8.45,4.38),(7.60,3.80),color=ABLATE,dashed=True,behind=main_panel)

    # Matrix
    add_text(s,9.62,2.03,0.44,0.25,"ID",size=8,bold=True,color=MUTED)
    add_text(s,10.10,2.03,1.62,0.25,"Change",size=8,bold=True,color=MUTED)
    add_text(s,11.78,2.03,0.92,0.25,"Measures",size=8,bold=True,color=MUTED)
    rows=[
        ("B0","Full structure","reference",OUTPUT,OUTPUT_BG),
        ("B1","No physical branch","dual graph",ABLATE,ABLATE_BG),
        ("B2","No SE","channel gate",RESIDUAL,RESIDUAL_BG),
        ("B3","No coord residual","coordinate",ABLATE,ABLATE_BG),
        ("B4","mean aggregation","aggregator",FEATURE,FEATURE_BG),
        ("B5","Include self in KNN","root repeat",PHYSICAL,PHYSICAL_BG),
        ("B6","No feature residual","residual",RESIDUAL,RESIDUAL_BG),
        ("B7","No explicit coord paths","interaction",ABLATE,ABLATE_BG),
        ("B8","EdgeCNN replaces SAGE","operator",FEATURE,FEATURE_BG),
    ]
    y=2.30
    for rid,change,measure,c,bg in rows:
        add_shape(s,9.62,y,0.46,0.42,rid,fill=bg,line=c,size=8.3,bold=True)
        add_shape(s,10.15,y,1.54,0.42,change,fill=WHITE,line=BORDER,accent=c,size=7.4)
        add_shape(s,11.76,y,0.94,0.42,measure,fill=WHITE,line=BORDER,size=7.2)
        y+=0.47
    add_text(s,9.62,6.63,3.05,0.18,"B0–B7: structural · B8: parameter-matched operator",size=7.2,bold=True,color=NAVY)
    return s


B_PARAMS = {
    "B0": 1331745,
    "B1": 1236001,
    "B2": 1286689,
    "B3": 1231389,
    "B4": 1331745,
    "B5": 1331745,
    "B6": 1287841,
    "B7": 1135645,
    "B8": 1331745,
}


def slide_structural_ablation(
    prs,
    exp_id,
    chinese_title,
    question,
    *,
    use_physical=True,
    use_se=True,
    use_coord=True,
    aggregation="max",
    exclude_self=True,
    feature_residual=True,
    interpretation="",
):
    """为一个 B 系列配置生成独立结构页。"""
    subtitle = f"{exp_id} · 单因素结构消融 · MLP head · λ_obj=0"
    s = setup_slide(prs, f"{exp_id} · {chinese_title}", subtitle, "model/graph_res_GCN_ablation.py")
    pipeline(s, "Ablation Blocks ×4\n32→64→64→128→256", head_text="MLP head\ncenter B×3")
    main_panel = add_panel(s, 0.34, 1.68, 8.72, 5.28, "Effective block structure · red dashed = physically absent")
    add_panel(s, 9.30, 1.68, 3.68, 5.28, "Configuration and interpretation")

    physical_off = not use_physical
    se_off = not use_se
    coord_off = not use_coord
    feature_res_off = not feature_residual

    # Feature stream.
    add_shape(s, 0.62, 2.30, 0.72, 0.50, "f", fill=FEATURE_BG, line=FEATURE, accent=FEATURE, bold=True)
    feature_label = f"Feature KNN\nSAGEConv_f · {aggregation}"
    add_shape(s, 1.62, 2.18, 1.48, 0.74, feature_label, fill=FEATURE_BG, line=FEATURE, accent=FEATURE, size=8.9, bold=True)
    add_shape(
        s, 3.42, 2.18, 1.20, 0.74,
        "OFF\nSE gate" if se_off else "SE gate\nGAP→MLP→σ",
        fill=ABLATE_BG if se_off else RESIDUAL_BG,
        line=ABLATE if se_off else RESIDUAL,
        accent=ABLATE if se_off else RESIDUAL,
        size=9.0, bold=True, dashed=se_off,
    )

    # Physical stream.
    add_shape(s, 0.62, 4.25, 0.72, 0.50, "p", fill=PHYSICAL_BG, line=PHYSICAL, accent=PHYSICAL, bold=True)
    physical_label = (
        "OFF\nCoordinate KNN + SAGEConv_p"
        if physical_off
        else f"Coordinate KNN\nSAGEConv_p · {aggregation}"
    )
    add_shape(
        s, 1.62, 4.10, 2.02, 0.80, physical_label,
        fill=ABLATE_BG if physical_off else PHYSICAL_BG,
        line=ABLATE if physical_off else PHYSICAL,
        accent=ABLATE if physical_off else PHYSICAL,
        size=8.8, bold=True, dashed=physical_off,
    )

    # Fusion and output.
    add_shape(s, 4.98, 3.02, 0.90, 0.76, "Feature only" if physical_off else "Cat", fill=FUSION_BG, line=FUSION, accent=FUSION, size=8.9, bold=True)
    add_shape(s, 6.18, 3.02, 0.76, 0.76, "Fuse", fill=FUSION_BG, line=FUSION, accent=FUSION, size=9.5, bold=True)
    add_shape(s, 7.24, 3.02, 0.88, 0.76, "Σ + Act", fill=OUTPUT_BG, line=OUTPUT, accent=OUTPUT, size=9.2, bold=True)
    add_shape(s, 8.36, 3.02, 0.50, 0.76, "f_out", fill=OUTPUT_BG, line=OUTPUT, size=8.7, bold=True)

    # Residual paths.
    add_shape(
        s, 1.72, 5.54, 1.80, 0.64,
        "OFF\nfeature residual" if feature_res_off else "feature residual\nIdentity / projection",
        fill=ABLATE_BG if feature_res_off else RESIDUAL_BG,
        line=ABLATE if feature_res_off else RESIDUAL,
        accent=ABLATE if feature_res_off else RESIDUAL,
        size=8.6, bold=feature_res_off, dashed=feature_res_off,
    )
    add_shape(
        s, 4.00, 5.42, 2.20, 0.84,
        "OFF\ncoordinate residual group" if coord_off else "coordinate residual group\ncoord_gate + coord_res + encoder",
        fill=ABLATE_BG if coord_off else FUSION_BG,
        line=ABLATE if coord_off else FUSION,
        accent=ABLATE if coord_off else FUSION,
        size=8.6, bold=coord_off, dashed=coord_off,
    )
    add_shape(
        s, 6.62, 5.54, 1.48, 0.64,
        "absent" if coord_off else "coord_scale × delta",
        fill=ABLATE_BG if coord_off else FUSION_BG,
        line=ABLATE if coord_off else FUSION,
        accent=ABLATE if coord_off else FUSION,
        size=8.5, dashed=coord_off,
    )

    # Main active/removed paths.
    route(s, [(1.34, 2.55), (1.62, 2.55)], color=FEATURE)
    if use_se:
        route(s, [(3.10, 2.55), (3.42, 2.55)], color=FEATURE)
        bezier(s, (4.62, 2.55), (4.78, 2.72), (4.78, 3.20), (4.98, 3.27), color=RESIDUAL, behind=main_panel)
    else:
        bezier(s, (3.10, 2.55), (3.80, 2.75), (4.50, 3.08), (4.98, 3.27), color=FEATURE, behind=main_panel)
        route(s, [(3.10, 2.70), (3.42, 2.70)], color=ABLATE, dashed=True)

    route(s, [(1.34, 4.50), (1.62, 4.50)], color=PHYSICAL if use_physical else ABLATE, dashed=physical_off)
    if use_physical:
        bezier(s, (3.64, 4.50), (4.18, 4.30), (4.55, 3.70), (4.98, 3.55), color=PHYSICAL, behind=main_panel)
    else:
        bezier(s, (3.64, 4.50), (4.10, 4.20), (4.55, 3.80), (4.98, 3.55), color=ABLATE, dashed=True, behind=main_panel)

    route(s, [(5.88, 3.40), (6.18, 3.40)], color=FUSION)
    route(s, [(6.94, 3.40), (7.24, 3.40)], color=OUTPUT)
    route(s, [(8.12, 3.40), (8.36, 3.40)], color=OUTPUT)

    if feature_residual:
        bezier(s, (0.98, 2.80), (0.86, 4.08), (1.05, 5.55), (1.72, 5.86), color=RESIDUAL, behind=main_panel)
        bezier(s, (3.52, 5.86), (4.42, 5.68), (5.86, 4.08), (7.24, 3.66), color=RESIDUAL, behind=main_panel)
    else:
        bezier(s, (0.98, 2.80), (0.86, 4.08), (1.05, 5.55), (1.72, 5.86), color=ABLATE, dashed=True, behind=main_panel)

    if use_coord:
        bezier(s, (0.98, 4.75), (1.72, 5.45), (3.80, 5.05), (5.10, 5.42), color=FUSION, behind=main_panel)
        bezier(s, (6.20, 5.86), (6.34, 5.82), (6.48, 5.82), (6.62, 5.86), color=FUSION, behind=main_panel)
        bezier(s, (8.10, 5.86), (8.62, 5.28), (8.48, 4.28), (7.68, 3.78), color=FUSION, behind=main_panel)
    else:
        bezier(s, (0.98, 4.75), (1.72, 5.45), (3.80, 5.05), (5.10, 5.42), color=ABLATE, dashed=True, behind=main_panel)

    # B5 explicitly visualizes self inclusion.
    if not exclude_self:
        pill(s, 1.78, 3.27, 1.26, "KNN includes i", PHYSICAL, PHYSICAL_BG)
        bezier(s, (2.02, 2.18), (1.76, 1.84), (2.96, 1.84), (2.70, 2.18), color=FEATURE, behind=main_panel)
        bezier(s, (2.04, 4.10), (1.74, 3.78), (3.30, 3.78), (3.00, 4.10), color=PHYSICAL, behind=main_panel)

    if aggregation == "mean":
        pill(s, 2.05, 3.27, 1.30, "AGG = mean", FEATURE, FEATURE_BG)

    # Right configuration cards.
    changed = []
    if physical_off: changed.append("physical branch = OFF")
    if se_off: changed.append("SE gate = OFF")
    if coord_off: changed.append("coord residual = OFF")
    if aggregation != "max": changed.append(f"aggregation = {aggregation}")
    if not exclude_self: changed.append("exclude_self = false")
    if feature_res_off: changed.append("feature residual = OFF")
    add_shape(s, 9.58, 2.08, 3.10, 0.70, "唯一有效变化\n" + " · ".join(changed), fill=ABLATE_BG if "OFF" in " ".join(changed) else FEATURE_BG, line=ABLATE if "OFF" in " ".join(changed) else FEATURE, accent=ABLATE if "OFF" in " ".join(changed) else FEATURE, size=9.2, bold=True)
    add_shape(s, 9.58, 2.98, 3.10, 0.82, f"研究问题\n{question}", fill=NEUTRAL_BG, line=NAVY, accent=NAVY, size=8.9, bold=True)
    config_text = (
        f"physical={str(use_physical).lower()} · SE={str(use_se).lower()}\n"
        f"coord_res={str(use_coord).lower()} · feature_res={str(feature_residual).lower()}\n"
        f"AGG={aggregation} · exclude_self={str(exclude_self).lower()}"
    )
    add_shape(s, 9.58, 4.02, 3.10, 0.88, config_text, fill=WHITE, line=BORDER, accent=PHYSICAL, size=8.5)
    add_shape(s, 9.58, 5.12, 3.10, 0.68, f"MLP head · effective λ_obj=0\nParameters: {B_PARAMS[exp_id]:,}", fill=OUTPUT_BG, line=OUTPUT, accent=OUTPUT, size=8.9, bold=True)
    add_shape(s, 9.58, 6.02, 3.10, 0.64, interpretation, fill=FUSION_BG, line=FUSION, accent=FUSION, size=8.2)
    return s


def slide_operator_comparison(prs):
    """B8: parameter-matched GraphSAGE and EdgeCNN operator control."""
    s = setup_slide(
        prs,
        "B8 \u00b7 GraphSAGE vs EdgeCNN",
        "\u53c2\u6570\u5339\u914d\u7b97\u5b50\u5bf9\u7167 \u00b7 \u4ec5\u6539\u53d8\u5c40\u90e8\u6d88\u606f\u4f20\u9012\u89c4\u5219",
        "model/graph_res_GCN.py \u00b7 model/graph_res_GCN_ablation.py",
    )
    pipeline(
        s,
        "Operator-controlled Blocks \u00d74\n32\u219264\u219264\u2192128\u2192256",
        head_text="MLP head\ncenter B\u00d73",
    )
    main_panel = add_panel(
        s, 0.34, 1.68, 8.86, 5.28,
        "Same KNN graph \u00b7 two parameter-matched local operators",
    )
    add_panel(s, 9.44, 1.68, 3.54, 5.28, "Controlled comparison")

    add_shape(s, 0.66, 2.17, 1.08, 0.62, "x \u00b7 B\u00d7Cin\u00d7N", fill=NEUTRAL_BG, line=NAVY, accent=NAVY, size=9.1, bold=True)
    add_shape(s, 2.02, 2.17, 1.34, 0.62, "Same KNN\nedge_index j\u2192i", fill=PHYSICAL_BG, line=PHYSICAL, accent=PHYSICAL, size=8.9, bold=True)
    route(s, [(1.74, 2.48), (2.02, 2.48)], color=NAVY)

    add_shape(s, 0.66, 3.22, 1.18, 0.58, "B0 \u00b7 GraphSAGE", fill=FEATURE_BG, line=FEATURE, accent=FEATURE, size=9.2, bold=True)
    add_shape(s, 2.18, 3.12, 1.42, 0.78, "AGG{x_j}\nmax / mean", fill=FEATURE_BG, line=FEATURE, accent=FEATURE, size=8.9)
    add_shape(s, 3.96, 3.12, 1.64, 0.78, "W_neighbor\u00b7AGG{x_j}\n+ W_root\u00b7x_i", fill=FEATURE_BG, line=FEATURE, accent=FEATURE, size=8.6, bold=True)
    add_shape(s, 5.96, 3.12, 1.06, 0.78, "h_i (SAGE)", fill=OUTPUT_BG, line=OUTPUT, accent=OUTPUT, size=9.4, bold=True)
    route(s, [(1.84, 3.51), (2.18, 3.51)], color=FEATURE)
    route(s, [(3.60, 3.51), (3.96, 3.51)], color=FEATURE)
    route(s, [(5.60, 3.51), (5.96, 3.51)], color=FEATURE)

    add_shape(s, 0.66, 4.54, 1.18, 0.58, "B8 \u00b7 EdgeCNN", fill=FUSION_BG, line=FUSION, accent=FUSION, size=9.2, bold=True)
    add_shape(s, 2.18, 4.42, 1.42, 0.82, "edge feature\n[x_j-x_i, x_i]", fill=FUSION_BG, line=FUSION, accent=FUSION, size=8.8)
    add_shape(s, 3.96, 4.42, 1.64, 0.82, "shared 1\u00d71 CNN\nLinear(2Cin,Cout)", fill=FUSION_BG, line=FUSION, accent=FUSION, size=8.7, bold=True)
    add_shape(s, 5.96, 4.42, 1.06, 0.82, "AGG\nh_i (CNN)", fill=OUTPUT_BG, line=OUTPUT, accent=OUTPUT, size=9.1, bold=True)
    route(s, [(1.84, 4.83), (2.18, 4.83)], color=FUSION)
    route(s, [(3.60, 4.83), (3.96, 4.83)], color=FUSION)
    route(s, [(5.60, 4.83), (5.96, 4.83)], color=FUSION)

    add_shape(s, 7.42, 3.62, 1.40, 0.86, "Same downstream\nBN + Act \u2192 SE\nDual fusion + residuals", fill=RESIDUAL_BG, line=RESIDUAL, accent=RESIDUAL, size=8.5, bold=True)
    bezier(s, (7.02, 3.51), (7.18, 3.51), (7.22, 3.82), (7.42, 3.92), color=FEATURE, behind=main_panel)
    bezier(s, (7.02, 4.83), (7.18, 4.83), (7.22, 4.33), (7.42, 4.18), color=FUSION, behind=main_panel)
    add_shape(s, 1.18, 5.78, 6.96, 0.66, "Exact parameter match: GraphSAGE = 1,331,745 \u00b7 EdgeCNN = 1,331,745", fill=OUTPUT_BG, line=OUTPUT, accent=OUTPUT, size=10.2, bold=True)
    add_text(s, 1.18, 6.50, 6.96, 0.20, "Compare B8-B0 per seed: classification, center error, mIoU and AP metrics.", size=8.2, bold=True, color=NAVY)

    add_shape(s, 9.72, 2.06, 2.98, 0.70, "\u552f\u4e00\u6709\u6548\u53d8\u5316\noperator: sage \u2192 edge_cnn", fill=FEATURE_BG, line=FEATURE, accent=FEATURE, size=9.2, bold=True)
    add_shape(s, 9.72, 2.96, 2.98, 1.04, "\u5b8c\u5168\u56fa\u5b9a\nKNN \u00b7 \u53cc\u5206\u652f \u00b7 AGG \u00b7 SE \u00b7 fusion\nfeature/coord residual \u00b7 MLP head", fill=NEUTRAL_BG, line=NAVY, accent=NAVY, size=8.5, bold=True)
    add_shape(s, 9.72, 4.20, 2.98, 0.76, "\u7814\u7a76\u95ee\u9898\n\u56fe\u5377\u79ef\u662f\u5426\u4f18\u4e8e CNN \u5f0f\u8fb9\u5377\u79ef\uff1f", fill=FUSION_BG, line=FUSION, accent=FUSION, size=9.0, bold=True)
    add_shape(s, 9.72, 5.16, 2.98, 0.70, "Registered seeds: 42 / 43\nFamily: operator", fill=PHYSICAL_BG, line=PHYSICAL, accent=PHYSICAL, size=9.0, bold=True)
    add_shape(s, 9.72, 6.04, 2.98, 0.62, "Analysis: scripts/analyze_gcn_vs_edge_cnn.py", fill=OUTPUT_BG, line=OUTPUT, accent=OUTPUT, size=8.2, bold=True)
    return s


def slide_lambda_sensitivity(prs):
    """C 系列目标性损失敏感性页，不虚构任何结果。"""
    s = setup_slide(prs, "C series · λ_obj sensitivity", "相同 centroid 架构，仅改变辅助目标性 BCE 权重", "utils/loss.py · scripts/train.py")
    pipeline(s, "GraphResidual-GCN ×4\ncentroid head enabled", head_text="Centroid head\ncenter B×3 + seg logits")
    main_panel = add_panel(s, 0.34, 1.68, 8.55, 5.28, "Centroid head and auxiliary objectness supervision")
    add_panel(s, 9.12, 1.68, 3.86, 5.28, "Sensitivity configurations")

    add_shape(s, 0.70, 2.34, 1.34, 0.70, "Point features\nB×512×N", fill=FEATURE_BG, line=FEATURE, accent=FEATURE, size=9.2, bold=True)
    add_shape(s, 2.38, 2.34, 1.34, 0.70, "seg_mlp\n512→128→1", fill=RESIDUAL_BG, line=RESIDUAL, accent=RESIDUAL, size=9.2, bold=True)
    add_shape(s, 4.06, 2.34, 1.16, 0.70, "seg_logits\nB×N", fill=RESIDUAL_BG, line=RESIDUAL, accent=RESIDUAL, size=9.2)
    add_shape(s, 5.56, 2.34, 1.22, 0.70, "softmax(s/τ)\nΣw=1", fill=FUSION_BG, line=FUSION, accent=FUSION, size=9.0)
    add_shape(s, 7.10, 2.34, 1.34, 0.70, "Σ w_i · xyz_i\ncenter B×3", fill=OUTPUT_BG, line=OUTPUT, accent=OUTPUT, size=9.1, bold=True)
    for x1, x2, color in [(2.04,2.38,FEATURE),(3.72,4.06,RESIDUAL),(5.22,5.56,RESIDUAL),(6.78,7.10,FUSION)]:
        route(s, [(x1,2.69),(x2,2.69)], color=color)

    add_shape(s, 0.70, 4.35, 1.34, 0.68, "GT box\nB×6", fill=PHYSICAL_BG, line=PHYSICAL, accent=PHYSICAL, size=9.2, bold=True)
    add_shape(s, 2.38, 4.35, 1.52, 0.68, "Point-in-box labels\n0 / 1", fill=PHYSICAL_BG, line=PHYSICAL, accent=PHYSICAL, size=8.8)
    add_shape(s, 4.24, 4.35, 1.22, 0.68, "BCE\nL_obj", fill=ABLATE_BG, line=ABLATE, accent=ABLATE, size=9.4, bold=True)
    add_shape(s, 5.82, 4.35, 1.24, 0.68, "λ_obj × L_obj", fill=FUSION_BG, line=FUSION, accent=FUSION, size=9.2, bold=True)
    add_shape(s, 7.40, 4.35, 1.04, 0.68, "L_total", fill=OUTPUT_BG, line=OUTPUT, accent=OUTPUT, size=9.4, bold=True)
    for x1, x2, color in [(2.04,2.38,PHYSICAL),(3.90,4.24,PHYSICAL),(5.46,5.82,ABLATE),(7.06,7.40,FUSION)]:
        route(s, [(x1,4.69),(x2,4.69)], color=color)
    bezier(s, (4.64,3.04), (4.72,3.55), (4.72,4.02), (4.85,4.35), color=ABLATE, behind=main_panel)
    add_text(s, 0.72, 5.52, 7.55, 0.52, "L_total = 1.0·L_cls + 10.0·L_depth + λ_obj·L_obj", size=14, bold=True, color=NAVY)
    add_text(s, 0.72, 6.18, 7.55, 0.32, "Architecture, split, seed protocol and initialization rules remain fixed; every λ run is trained from scratch.", size=8.8, color=MUTED)

    lambdas = [("0", "A2 · reuse", OUTPUT), ("0.25", "low · new train", PHYSICAL), ("0.5", "A3 · reuse", FUSION), ("1.0", "strong · new train", ABLATE)]
    y = 2.15
    for value, label, color in lambdas:
        add_shape(s, 9.42, y, 0.72, 0.55, value, fill=WHITE, line=color, accent=color, size=10, bold=True)
        add_shape(s, 10.28, y, 2.38, 0.55, label, fill=NEUTRAL_BG, line=BORDER, accent=color, size=8.8)
        y += 0.82
    add_shape(s, 9.42, 5.78, 3.24, 0.80, "Frozen sensitivity grid\n{0, 0.25, 0.5, 1.0}\n0 / 0.5 reuse A2 / A3", fill=FUSION_BG, line=FUSION, accent=FUSION, size=9.0, bold=True)
    return s

def new_presentation():
    p=Presentation(); p.slide_width=Inches(W); p.slide_height=Inches(H)
    p.core_properties.title="SPAD editable model architecture"
    p.core_properties.subject="GraphResidual, formal GraphResidual-GCN, structural ablations, and parameter-matched EdgeCNN control"
    p.core_properties.author="SPAD Project"
    p.core_properties.comments="Generated by scripts/generate_model_architecture_ppt.py after local render-review iterations."
    return p


def save_one(path:Path,builder:Callable):
    p=new_presentation(); builder(p); path.parent.mkdir(parents=True,exist_ok=True); p.save(path)


def run_with_config(cfg:ScriptConfig):
    out=cfg.output_dir.resolve(); out.mkdir(parents=True,exist_ok=True)
    combined_path=out/"SPAD_model_architectures_editable_combined.pptx"

    # The user requested one consolidated deck only; remove obsolete one-slide decks.
    for obsolete in (
        "01_GraphResidual_EdgeConv_original_editable.pptx",
        "02_GraphResidual_GCN_formal_editable.pptx",
        "03_GraphResidual_GCN_ablation_editable.pptx",
    ):
        path=out/obsolete
        if path.is_file(): path.unlink()

    combined=new_presentation()
    slide_edgeconv(combined)
    slide_gcn(combined)
    slide_ablation(combined)
    slide_structural_ablation(
        combined,"B1","关闭坐标图分支","双图中的显式坐标 GraphSAGE 分支是否必要？",
        use_physical=False,
        interpretation="仅移除 physical branch；SE、feature residual 与 coordinate residual 保留。",
    )
    slide_structural_ablation(
        combined,"B2","关闭 SE","通道重标定是否提供独立收益？",
        use_se=False,
        interpretation="特征 SAGE 输出直接进入双流融合；其余结构不变。",
    )
    slide_structural_ablation(
        combined,"B3","关闭坐标残差","显式坐标残差注入是否必要？",
        use_coord=False,
        interpretation="移除 coord_gate、coord_res、coord_encoder 与 coord_scale。",
    )
    slide_structural_ablation(
        combined,"B4","max 改为 mean 聚合","稀疏目标是否更适合 max 而不是 mean？",
        aggregation="mean",
        interpretation="Feature 与 physical 两个 SAGEConv 均使用 mean 聚合。",
    )
    slide_structural_ablation(
        combined,"B5","KNN 包含自身点","root transform 外再次聚合自身是否有影响？",
        exclude_self=False,
        interpretation="KNN 保留中心点；SAGEConv 的 root transform 仍然存在。",
    )
    slide_structural_ablation(
        combined,"B6","关闭特征残差","显式 feature residual 是否改善优化与语义保留？",
        feature_residual=False,
        interpretation="仅移除 block 的 Identity / projection feature skip。",
    )
    slide_structural_ablation(
        combined,"B7","关闭两条坐标路径","physical branch 与 coordinate residual 是否存在交互？",
        use_physical=False,use_coord=False,
        interpretation="Stem 仍输入 xyzi；B7 不是无坐标模型，而是移除两条显式坐标增强路径。",
    )
    slide_operator_comparison(combined)
    slide_lambda_sensitivity(combined)
    combined.save(combined_path)
    return [combined_path]

def build_parser():
    p=argparse.ArgumentParser(description="Generate polished editable SPAD model PPTX")
    p.add_argument("--output-dir",type=Path,default=Path("model/ppt"),help="Output directory")
    return p


def main(argv:Optional[Sequence[str]]=None):
    args=build_parser().parse_args(argv)
    for f in run_with_config(ScriptConfig(args.output_dir)): print(f)
    return 0


def main_without_cli():
    for f in run_with_config(ScriptConfig(Path("model/ppt"))): print(f)


if __name__=="__main__":
    import sys
    raise SystemExit(main() if len(sys.argv)>1 else (main_without_cli() or 0))
